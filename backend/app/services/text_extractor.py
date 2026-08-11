import io
import logging
from typing import List, Tuple, Optional
import tiktoken

logger = logging.getLogger(__name__)


class TextExtractor:
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        try:
            self.encoding = tiktoken.get_encoding("cl100k_base")
        except Exception:
            self.encoding = None
    
    def extract_text(self, file_data: bytes, file_type: str) -> Optional[str]:
        try:
            if file_type.lower() == "pdf":
                return self._extract_from_pdf(file_data)
            elif file_type.lower() in ["docx", "doc"]:
                return self._extract_from_docx(file_data)
            elif file_type.lower() in ["xlsx", "xls"]:
                return self._extract_from_xlsx(file_data)
            elif file_type.lower() in ["pptx", "ppt"]:
                return self._extract_from_pptx(file_data)
            elif file_type.lower() in ["png", "jpg", "jpeg", "gif", "webp", "bmp"]:
                return self._extract_from_image(file_data, file_type)
            elif file_type.lower() == "txt":
                return file_data.decode("utf-8", errors="ignore")
            else:
                logger.warning(f"Unsupported file type: {file_type}")
                return None
        except Exception as e:
            logger.error(f"Error extracting text from {file_type}: {e}")
            return None
    
    def _extract_from_pdf(self, file_data: bytes) -> str:
        try:
            import fitz
            doc = fitz.open(stream=file_data, filetype="pdf")
            text_parts = []
            for page in doc:
                text = page.get_text()  # type: ignore[attr-defined]
                if text:
                    text_parts.append(text)
            doc.close()
            return "\n\n".join(text_parts)
        except ImportError:
            logger.error("PyMuPDF (fitz) not installed. Run: pip install pymupdf")
            return ""
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            return ""
    
    def _extract_from_docx(self, file_data: bytes) -> str:
        try:
            from docx import Document
            doc = Document(io.BytesIO(file_data))
            text_parts = []
            for para in doc.paragraphs:
                if para.text:
                    text_parts.append(para.text)
            return "\n\n".join(text_parts)
        except ImportError:
            logger.error("python-docx not installed")
            return ""
    
    def _extract_from_xlsx(self, file_data: bytes) -> str:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_data), data_only=True)
            text_parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell else "" for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
            return "\n".join(text_parts)
        except ImportError:
            logger.error("openpyxl not installed")
            return ""
    
    def _extract_from_pptx(self, file_data: bytes) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_data))
            text_parts = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:  # type: ignore[attr-defined]
                        for para in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                            if para.text.strip():
                                slide_texts.append(para.text)
                    if shape.has_table:  # type: ignore[attr-defined]
                        for row in shape.table.rows:  # type: ignore[attr-defined]
                            row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                            if row_text:
                                slide_texts.append(row_text)
                if slide_texts:
                    text_parts.append(f"[第{slide_idx}页]\n" + "\n".join(slide_texts))
            return "\n\n".join(text_parts)
        except ImportError:
            logger.error("python-pptx not installed. Run: pip install python-pptx")
            return ""
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {e}")
            return ""
    
    def _extract_from_image(self, file_data: bytes, file_type: str) -> str:
        logger.warning("Image OCR requires async OCR service, use ocr_service.extract_text_from_image() instead")
        return ""
    
    def count_tokens(self, text: str) -> int:
        if self.encoding:
            return len(self.encoding.encode(text))
        return len(text.split())
    
    def split_text_into_chunks(self, text: str) -> List[Tuple[str, int]]:
        if not text:
            return []
        
        if self.encoding:
            tokens = self.encoding.encode(text)
            chunks = []
            start = 0
            chunk_index = 0
            
            while start < len(tokens):
                end = min(start + self.chunk_size, len(tokens))
                chunk_tokens = tokens[start:end]
                chunk_text = self.encoding.decode(chunk_tokens)
                chunks.append((chunk_text, chunk_index))
                chunk_index += 1
                
                if end >= len(tokens):
                    break
                
                start = end - self.chunk_overlap
            
            return chunks
        else:
            words = text.split()
            chunks = []
            start = 0
            chunk_index = 0
            
            while start < len(words):
                end = min(start + self.chunk_size, len(words))
                chunk_text = " ".join(words[start:end])
                chunks.append((chunk_text, chunk_index))
                chunk_index += 1
                
                if end >= len(words):
                    break
                
                start = end - self.chunk_overlap
            
            return chunks
